import os
import subprocess
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def convert_word_to_ppt(input_path, output_dir):
    subprocess.run([
        'soffice', '--headless', 
        '--convert-to', 'pptx', 
        '--outdir', output_dir, 
        input_path
    ], check=True)
    filename = os.path.basename(input_path).rsplit('.', 1)[0] + '.pptx'
    return os.path.join(output_dir, filename)

def convert_ppt_to_word(input_path, output_dir):
    subprocess.run([
        'soffice', '--headless', 
        '--convert-to', 'docx', 
        '--outdir', output_dir, 
        input_path
    ], check=True)
    filename = os.path.basename(input_path).rsplit('.', 1)[0] + '.docx'
    return os.path.join(output_dir, filename)

def convert_pdf_to_ppt(input_path, output_path):
    images = convert_from_path(input_path)
    prs = Presentation()
    
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    
    for i, img in enumerate(images):
        temp_img_path = f"/tmp/temp_slide_{i}.png"
        img.save(temp_img_path, "PNG")
        
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        if os.path.exists(temp_img_path):
            os.remove(temp_img_path)
            
    prs.save(output_path)